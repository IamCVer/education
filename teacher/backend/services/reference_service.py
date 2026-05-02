from __future__ import annotations

import mimetypes
import uuid
from pathlib import Path
from typing import Dict, Tuple


TEXT_LIMIT = 2000


def detect_file_type(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix == ".docx":
        return "docx"
    if suffix == ".pptx":
        return "pptx"
    if suffix in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}:
        return "image"
    if suffix in {".mp4", ".mov", ".avi", ".mkv", ".webm"}:
        return "video"
    if suffix in {".txt", ".md"}:
        return "text"
    return "file"


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except Exception:
        return ""
    try:
        reader = PdfReader(str(path))
        texts = [(page.extract_text() or "") for page in reader.pages[:5]]
        return "\n".join(texts)
    except Exception:
        return ""


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except Exception:
        return ""
    try:
        doc = Document(str(path))
        return "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
    except Exception:
        return ""


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return ""
    try:
        presentation = Presentation(str(path))
        texts = []
        for slide in presentation.slides[:8]:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    texts.append(text)
        return "\n".join(texts)
    except Exception:
        return ""


def extract_reference_text(path: Path, file_type: str) -> str:
    if file_type == "pdf":
        return _read_pdf(path)
    if file_type == "docx":
        return _read_docx(path)
    if file_type == "pptx":
        return _read_pptx(path)
    if file_type == "text":
        return path.read_text(encoding="utf-8", errors="ignore")
    return ""


def summarize_reference(path: Path, purpose: str = "") -> Tuple[str, str]:
    file_type = detect_file_type(path)
    preview_text = extract_reference_text(path, file_type).strip()
    if preview_text:
        preview_text = preview_text[:TEXT_LIMIT]
        summary = preview_text[:240].replace("\n", " ")
        return summary, preview_text
    mime = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    summary = f"{file_type} 文件，MIME={mime}"
    if purpose:
        summary = f"{summary}；用途：{purpose}"
    return summary, ""


async def save_reference_file(upload_file, session_dir: Path, purpose: str = "") -> Dict[str, str]:
    file_id = str(uuid.uuid4())
    original_name = upload_file.filename or "unnamed-file"
    safe_name = f"{file_id}-{original_name}"
    target = session_dir / safe_name
    content = await upload_file.read()
    target.write_bytes(content)
    file_type = detect_file_type(target)
    summary, preview_text = summarize_reference(target, purpose)
    return {
        "file_id": file_id,
        "file_name": original_name,
        "stored_name": safe_name,
        "file_type": file_type,
        "purpose": purpose,
        "summary": summary,
        "preview_text": preview_text,
        "parse_status": "parsed" if preview_text or file_type in {"image", "video"} else "saved",
        "path": str(target),
    }
